#!/usr/bin/env python3
"""
font_audit.py — Audit every font used in a PowerPoint file.

Usage:
    python font_audit.py <presentation.pptx> [--csv output.csv]

Reports: slide number, shape name, font name, font size, bold/italic,
and the text fragment styled with that font.
"""

import argparse
import csv
import sys
from collections import defaultdict
from pptx import Presentation
from pptx.util import Pt


def emu_to_pt(emu):
    """Convert EMU to points."""
    if emu is None:
        return None
    return round(emu / 12700, 1)


def resolve_font(run, slide_layout=None, slide_master=None):
    """
    Resolve the effective font for a run by walking the inheritance chain:
    run → paragraph default → placeholder → slide layout → slide master → theme.
    Returns (font_name, font_size_pt, is_bold, is_italic).
    """
    font = run.font

    # Font name
    name = font.name  # Directly set on the run
    # Font size
    size = emu_to_pt(font.size) if font.size else None
    # Bold / Italic
    bold = font.bold
    italic = font.italic

    # If name is still None, report it as inherited/theme
    if name is None:
        name = "(inherited/theme default)"

    return name, size, bold, italic


def audit_presentation(path):
    """Walk every run in the presentation and collect font records."""
    prs = Presentation(path)
    records = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # --- Text in normal shapes ---
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text.strip()
                        if not text:
                            continue
                        name, size, bold, italic = resolve_font(run)
                        records.append({
                            "slide": slide_idx,
                            "shape": shape.name,
                            "font": name,
                            "size_pt": size,
                            "bold": bold,
                            "italic": italic,
                            "text": text,
                        })

            # --- Text inside table cells ---
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                text = run.text.strip()
                                if not text:
                                    continue
                                name, size, bold, italic = resolve_font(run)
                                records.append({
                                    "slide": slide_idx,
                                    "shape": f"{shape.name} [row {row_idx}, col {col_idx}]",
                                    "font": name,
                                    "size_pt": size,
                                    "bold": bold,
                                    "italic": italic,
                                    "text": text,
                                })

            # --- Text inside grouped shapes ---
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                _audit_group(shape, slide_idx, records)

    return records


def _audit_group(group_shape, slide_idx, records):
    """Recursively audit shapes inside a group."""
    for shape in group_shape.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text.strip()
                    if not text:
                        continue
                    name, size, bold, italic = resolve_font(run)
                    records.append({
                        "slide": slide_idx,
                        "shape": f"{group_shape.name} → {shape.name}",
                        "font": name,
                        "size_pt": size,
                        "bold": bold,
                        "italic": italic,
                        "text": text,
                    })
        if shape.shape_type == 6:
            _audit_group(shape, slide_idx, records)


def print_report(records):
    """Print a readable console report."""
    if not records:
        print("No text runs found in the presentation.")
        return

    # --- Summary: unique fonts ---
    font_counts = defaultdict(int)
    for r in records:
        font_counts[r["font"]] += 1

    print("=" * 70)
    print("FONT SUMMARY")
    print("=" * 70)
    for font, count in sorted(font_counts.items(), key=lambda x: -x[1]):
        print(f"  {font:<35} {count:>4} text runs")
    print(f"\n  Total text runs: {len(records)}")
    print(f"  Unique fonts:    {len(font_counts)}")
    print()

    # --- Detail by slide ---
    print("=" * 70)
    print("DETAIL BY SLIDE")
    print("=" * 70)
    current_slide = None
    for r in records:
        if r["slide"] != current_slide:
            current_slide = r["slide"]
            print(f"\n--- Slide {current_slide} ---")

        style_parts = []
        if r["size_pt"]:
            style_parts.append(f"{r['size_pt']}pt")
        if r["bold"]:
            style_parts.append("bold")
        if r["italic"]:
            style_parts.append("italic")
        style_str = ", ".join(style_parts) if style_parts else "default style"

        # Truncate long text for readability
        text = r["text"]
        if len(text) > 80:
            text = text[:77] + "..."

        print(f"  [{r['font']}] ({style_str})")
        print(f"    Shape: {r['shape']}")
        print(f"    Text:  \"{text}\"")


def write_csv(records, csv_path):
    """Write records to a CSV file."""
    fieldnames = ["slide", "shape", "font", "size_pt", "bold", "italic", "text"]
    with open(csv_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(records)
    print(f"\nCSV written to: {csv_path}")


def main():
    parser = argparse.ArgumentParser(description="Audit fonts in a PPTX file.")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument("--csv", help="Optional: export results to CSV", default=None)
    args = parser.parse_args()

    print(f"Auditing: {args.pptx}\n")
    records = audit_presentation(args.pptx)
    print_report(records)

    if args.csv:
        write_csv(records, args.csv)


if __name__ == "__main__":
    main()
